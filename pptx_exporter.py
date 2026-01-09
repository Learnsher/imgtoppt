# ---- PPTX Exporter (add to your_converter.py) ----
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

class PPTXExporter:
    """
    Create PPTX slides with:
    - background image (cleaned)
    - editable textboxes overlaying OCR regions
    """

    def __init__(self, slide_width_in=13.333, slide_height_in=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width_in)
        self.prs.slide_height = Inches(slide_height_in)

    def add_slide_with_overlay(self, bg_image_path, text_regions):
        blank_layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(blank_layout)

        # background image fill
        img = Image.open(bg_image_path)
        img_w, img_h = img.size

        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        scale = max(slide_w / img_w, slide_h / img_h)
        pic_w = int(img_w * scale)
        pic_h = int(img_h * scale)
        left = int((slide_w - pic_w) / 2)
        top = int((slide_h - pic_h) / 2)

        slide.shapes.add_picture(bg_image_path, left, top, pic_w, pic_h)

        # text overlay
        scale_x = slide_w / img_w
        scale_y = slide_h / img_h

        for r in text_regions:
            x, y, w, h = r["bbox"]
            text = r["text"]

            tb = slide.shapes.add_textbox(
                int(x * scale_x),
                int(y * scale_y),
                max(1, int(w * scale_x)),
                max(1, int(h * scale_y)),
            )
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text

            # font size estimate
            font_size = max(8, int(h * scale_y * 0.7))
            p.font.size = Pt(font_size)

            # visually hide box
            tb.line.fill.background()
            tb.fill.background()

    def save(self, output_path):
        self.prs.save(output_path)
