# your_converter.py - 完整版本

from dataclasses import dataclass
from typing import List, Tuple, Optional
import numpy as np
from PIL import Image, ImageDraw
from scipy import ndimage
import easyocr
from simple_lama_inpainting import SimpleLama


@dataclass
class TextRegion:
    bbox: Tuple[int, int, int, int]
    text: str
    confidence: float
    quad: Optional[List] = None


class EditableDocConverter:
    def __init__(
        self,
        lang: str = 'ch',
        use_gpu: bool = False,
        min_confidence: float = 0.50,
    ):
        # 初始化 EasyOCR
        self.ocr = easyocr.Reader(['ch_sim', 'en'], gpu=use_gpu)
        self.inpainter = SimpleLama()
        self.min_confidence = min_confidence

    def process_document(
        self,
        image_path: str,
        clean_image_path: Optional[str] = None,
        dilation_size: int = 5,
        dilation_iter: int = 2,
        return_mask: bool = False,
    ):
        """
        處理文檔：OCR → 創建 mask → inpainting
        """
        image_pil = Image.open(image_path).convert("RGB")
        img_np = np.array(image_pil)

        # OCR 識別
        regions = self._extract_text_regions(img_np)

        # 創建 mask
        mask_pil = self._create_binary_mask_scipy(
            img_size=image_pil.size,
            regions=regions,
            dilation_size=dilation_size,
            dilation_iter=dilation_iter,
        )

        # Inpainting
        clean_pil = self.inpainter(image_pil, mask_pil)

        if clean_image_path:
            clean_pil.save(clean_image_path)

        # 轉換為可序列化格式
        text_regions = [
            {
                "bbox": r.bbox,
                "text": r.text,
                "confidence": r.confidence,
                "quad": r.quad,
            }
            for r in regions
        ]

        if return_mask:
            return clean_pil, text_regions, mask_pil
        return clean_pil, text_regions

    def _extract_text_regions(self, img_np) -> List[TextRegion]:
        """
        用 EasyOCR 識別文字
        返回格式：[[bbox, text, confidence], ...]
        """
        results = self.ocr.readtext(img_np)
        regions = []

        for (bbox, text, conf) in results:
            if conf < self.min_confidence or not text.strip():
                continue

            # bbox 係 4 點座標 [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
            xs = [p[0] for p in bbox]
            ys = [p[1] for p in bbox]
            x0, y0 = int(min(xs)), int(min(ys))
            w, h = int(max(xs) - x0), int(max(ys) - y0)

            regions.append(
                TextRegion(
                    bbox=(x0, y0, w, h),
                    text=text,
                    confidence=conf,
                    quad=bbox,
                )
            )
        return regions

    def _create_binary_mask_scipy(
        self,
        img_size: Tuple[int, int],
        regions: List[TextRegion],
        dilation_size: int = 5,
        dilation_iter: int = 2,
    ) -> Image.Image:
        """
        創建二值 mask（用 scipy 做 dilation）
        """
        w, h = img_size
        mask = Image.new("L", (w, h), 0)
        draw = ImageDraw.Draw(mask)

        # 繪製文字區域
        for r in regions:
            x, y, bw, bh = r.bbox
            draw.rectangle([x, y, x + bw, y + bh], fill=255)

        # 用 scipy 做 morphological dilation
        mask_np = np.array(mask, dtype=np.uint8)
        structure = np.ones((dilation_size, dilation_size), dtype=np.uint8)

        for _ in range(dilation_iter):
            mask_np = ndimage.binary_dilation(mask_np, structure=structure).astype(np.uint8) * 255

        return Image.fromarray(mask_np)


# ---- PPTX Exporter ----
from pptx import Presentation
from pptx.util import Inches, Pt


class PPTXExporter:
    """
    創建 PPTX slides：
    - 背景圖（cleaned）
    - 可編輯 textboxes（OCR 識別的位置）
    """

    def __init__(self, slide_width_in=13.333, slide_height_in=7.5):
        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width_in)
        self.prs.slide_height = Inches(slide_height_in)

    def add_slide_with_overlay(self, bg_image_path, text_regions):
        blank_layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(blank_layout)

        # 背景圖
        img = Image.open(bg_image_path)
        img_w, img_h = img.size
        slide_w = self.prs.slide_width
        slide_h = self.prs.slide_height

        scale = max(slide_w / img_w, slide_h / img_h)
        pic_w = int(img_w * scale)
        pic_h = int(img_h * scale)
        left = int((slide_w - pic_w) / 2)
        top = int((slide_h - pic_h) / 2)

        slide.shapes.add_picture(bg_image_path, left, top, pic_w, pic_h)

        # 文字圖層
        scale_x = slide_w / img_w
        scale_y = slide_h / img_h

        for r in text_regions:
            x, y, w, h = r["bbox"]
            text = r["text"]

            tb = slide.shapes.add_textbox(
                int(x * scale_x),
                int(y * scale_y),
                max(1, int(w * scale_x)),
                max(1, int(h * scale_y)),
            )
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text

            # 字體大小估算
            font_size = max(8, int(h * scale_y * 0.7))
            p.font.size = Pt(font_size)

            # 隱藏 textbox 邊框
            tb.line.fill.background()
            tb.fill.background()

    def save(self, output_path):
        self.prs.save(output_path)
